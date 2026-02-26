import json
from pptx import Presentation
import google.genai as genai
from google.genai import types
from django.conf import settings

def extract_ppt_text(file_path):
    """
    Extracts text from a PowerPoint file.
    """
    prs = Presentation(file_path)
    text_content = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text_content.append("\n".join(slide_text))

    return "\n\n".join(text_content)

def get_gemini_client():
    """
    Returns a configured Gemini API client.
    """
    api_key = getattr(settings, 'GEMINI_API_KEY', None)
    if not api_key or api_key == 'PLACEHOLDER_KEY':
        # Fallback or error handling - in production, raise proper error
        print("Warning: GEMINI_API_KEY not set.")
        return None
    return genai.Client(api_key=api_key)

def summarize_text(text):
    """
    Summarizes the given text using Google Gemini 2.5 Flash.
    Returns a JSON string with structured data.
    """
    client = get_gemini_client()
    if not client:
        return json.dumps({
            "summary": "AI Summarization unavailable: API Key missing.",
            "key_concepts": [],
            "important_terms": []
        })

    # JSON schema for the response
    schema = {
        "type": "OBJECT",
        "properties": {
            "summary": {"type": "STRING"},
            "key_concepts": {
                "type": "ARRAY",
                "items": {
                    "type": "OBJECT",
                    "properties": {
                        "title": {"type": "STRING"},
                        "description": {"type": "STRING"},
                        "color": {"type": "STRING", "enum": ["green", "blue", "purple"]}
                    },
                    "required": ["title", "description", "color"]
                }
            },
            "important_terms": {
                "type": "ARRAY",
                "items": {"type": "STRING"}
            }
        },
        "required": ["summary", "key_concepts", "important_terms"]
    }

    prompt = f"""
    Analyze the following presentation content and provide a structured learning summary.
    
    1. **summary**: A concise executive summary paragraph (approx 3-5 sentences) suitable for a quick overview.
    2. **key_concepts**: Extract 3-5 key concepts. For each, provide a short 'title', a 'description', and assign a 'color' (green, blue, or purple) based on category/importance.
    3. **important_terms**: List specific nouns/verbs that are crucial for sign language practice.
    
    Content:
    {text[:10000]}
    """
    
    try:
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type='application/json',
                response_schema=schema
            )
        )
        return response.text
    except Exception as e:
        print(f"Summary Generation Error: {e}")
        return json.dumps({
            "summary": f"Error generating summary: {str(e)}",
            "key_concepts": [],
            "important_terms": []
        })

def generate_mcq(text, num_questions=5, difficulty='Medium'):
    """
    Generates multiple-choice questions from the text using Google Gemini.
    Returns a list of dictionaries.
    """
    client = get_gemini_client()
    if not client:
        return []

    # JSON schema for the response
    schema = {
        "type": "ARRAY",
        "items": {
            "type": "OBJECT",
            "properties": {
                "question": {"type": "STRING"},
                "options": {
                    "type": "ARRAY",
                    "items": {"type": "STRING"}
                },
                "correct_answer": {"type": "STRING"},
                "explanation": {"type": "STRING"}
            },
            "required": ["question", "options", "correct_answer", "explanation"]
        }
    }

    prompt = f"""
    Generate {num_questions} multiple-choice questions from the following text.
    Difficulty Level: {difficulty}.
    
    The output must be a valid JSON array of objects.
    Each object must have:
    - question: The question text
    - options: An array of 4 distinct options
    - correct_answer: The correct option text (must match one of the options)
    - explanation: A brief explanation of why the answer is correct
    
    Text Content:
    {text[:15000]}
    """

    try:
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type='application/json',
                response_schema=schema
            )
        )
        return json.loads(response.text)
    except Exception as e:
        print(f"MCQ Generation Error: {e}")
        return []

import nltk
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer
from django.contrib.staticfiles import finders
import string

def process_text_for_sign_language(text):
    """
    Processes text to return a list of words suitable for sign language animation.
    Filters stop words, lemmatizes, and checks for file existence.
    """
    if not text:
        return []

    text = text.lower()
    try:
        words = word_tokenize(text)
    except LookupError:
        nltk.download('punkt')
        words = word_tokenize(text)

    try:
        tagged = nltk.pos_tag(words)
    except LookupError:
        nltk.download('averaged_perceptron_tagger')
        tagged = nltk.pos_tag(words)

    # Simplified tense logic (defaulting to present for summary)
    
    stop_words = set(["mightn't", 're', 'wasn', 'wouldn', 'be', 'has', 'that', 'does', 'shouldn', 'do', "you've",'off', 'for', "didn't", 'm', 'ain', 'haven', "weren't", 'are', "she's", "wasn't", 'its', "haven't", "wouldn't", 'don', 'weren', 's', "you'd", "don't", 'doesn', "hadn't", 'is', 'was', "that'll", "should've", 'a', 'then', 'the', 'mustn', 'i', 'nor', 'as', "it's", "needn't", 'd', 'am', 'have',  'hasn', 'o', "aren't", "you'll", "couldn't", "you're", "mustn't", 'didn', "doesn't", 'll', 'an', 'hadn', 'whom', 'y', "hasn't", 'itself', 'couldn', 'needn', "shan't", 'isn', 'been', 'such', 'shan', "shouldn't", 'aren', 'being', 'were', 'did', 'ma', 't', 'having', 'mightn', 've', "isn't", "won't"])

    lr = WordNetLemmatizer()
    try:
        lr.lemmatize('test')
    except LookupError:
        nltk.download('wordnet')
        nltk.download('omw-1.4')

    filtered_text = []
    for w, p in zip(words, tagged):
        if w not in stop_words and w not in string.punctuation:
            if p[1] in ['VBG', 'VBD', 'VBZ', 'VBN', 'NN']:
                filtered_text.append(lr.lemmatize(w, pos='v'))
            elif p[1] in ['JJ', 'JJR', 'JJS', 'RBR', 'RBS']:
                filtered_text.append(lr.lemmatize(w, pos='a'))
            else:
                filtered_text.append(lr.lemmatize(w))

    # Match with existing assets
    final_words = []
    for w in filtered_text:
        # Capitalize first letter as assets are mostly TitleCase or uppercase
        w_title = w.capitalize() 
        path = w_title + ".mp4"
        
        # Check specific naming conventions if needed (e.g. "Do Not")
        # For now, simplistic check
        
        if finders.find(path):
            final_words.append(w_title)
        else:
            # If word not found, spell it out (return individual chars)
            # Or just append the word and let frontend handle spelling if missing?
            # The existing logic split it. Let's split it here to be safe and consistent.
            for c in w:
                if c.isalnum():
                    final_words.append(c.upper())

    return final_words
